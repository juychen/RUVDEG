"""Build a minimal 5-slide PPTX for RUVVAE-DEG. Text + necessary figures only."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image
import os

SANS = "Noto Sans CJK SC"
MONO = "Noto Sans Mono CJK SC"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x73, 0x73, 0x73)
ACCENT = RGBColor(0xB0, 0x3A, 0x2E)
RULE = RGBColor(0xDC, 0xDC, 0xDC)

SW, SH = 13.3333, 7.5
ML = 0.62
CW = SW - 2 * ML
F = "slide_figs"


# ---------------------------------------------------------------- helpers
def set_ea(run, name):
    """python-pptx only sets the latin typeface; add ea/cs so CJK renders."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.makeelement(qn(tag), {"typeface": name})
        rPr.append(el)


def style_of(s):
    bold = "b" in s
    mono = "c" in s
    col = ACCENT if "a" in s else (GRAY if "g" in s else INK)
    return bold, mono, col


def tbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


_used = {}


def P(tf, spans, size=13, sb=0, sa=5, line=1.22):
    key = id(tf)
    if key not in _used:
        p, _used[key] = tf.paragraphs[0], True
    else:
        p = tf.add_paragraph()
    p.space_before, p.space_after, p.line_spacing = Pt(sb), Pt(sa), line
    if isinstance(spans, str):
        spans = [(spans, "")]
    for item in spans:
        text, sty = (item, "") if isinstance(item, str) else item
        bold, mono, col = style_of(sty)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = MONO if mono else SANS
        r.font.bold = bold
        r.font.color.rgb = col
        set_ea(r, r.font.name)
    return p


def rule(slide, x, y, w, color=RULE, h=0.014):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_pic(slide, name, x, y, w, h):
    path = os.path.join(F, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w / h > ar:
        ph, pw = h, h * ar
    else:
        pw, ph = w, w / ar
    slide.shapes.add_picture(path, Inches(x + (w - pw) / 2),
                             Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))


def new_slide(prs, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = tbox(s, ML, 0.42, CW, 0.5)
    P(tf, [(title, "b")], size=25, sa=0)
    if kicker:
        tf2 = tbox(s, ML, 0.95, CW, 0.3)
        P(tf2, [(kicker, "g")], size=11.5, sa=0)
        rule(s, ML, 1.32, CW)
        return s, 1.52
    rule(s, ML, 1.06, CW)
    return s, 1.26


# ---------------------------------------------------------------- KL crop
os.makedirs(F, exist_ok=True)
_l = Image.open("ppt_figs/fig1_training_loss.png")
_w, _h = _l.size
_l.crop((int(_w * 0.665), 0, _w, _h)).save(os.path.join(F, "kl.png"))

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)

# ================================================================ SLIDE 1
s, y = new_slide(
    prs, "RUVVAE-DEG · 问题与思路",
    "小鼠丘脑 TH Tll1_Thsd7b Glut 神经元 · CSRES (慢性社会挫败应激) vs CON · "
    "4 821 细胞 × 16 428 基因")

tf = tbox(s, ML, y, CW * 0.47, 5.4)
P(tf, [("核心困难：donor 与 status 完全共线", "b")], size=16, sa=14)
P(tf, [("CSRS1-3 / CSRS9-1 / CSRS10-3", "c"), ("  →  全部 CSRES", ""),
       ("  2 140 细胞", "g")], size=12.5, sa=6)
P(tf, [("MW22B / MW45A / MW47A / MW51A", "c"), ("  →  全部 CON", ""),
       ("  2 681 细胞", "g")], size=12.5, sa=14)
P(tf, [("送样公司 3 家 · 建库批次 4 个 · 测序深度 3 548 ± 1 202 基因/细胞", "g")],
  size=12, sa=16)
P(tf, [("任何「按组求均值相减」都无法区分应激效应与批次效应。", "ab")],
  size=13.5, sa=34)

P(tf, [("经典 RUV 的困境", "b")], size=16, sa=14)
P(tf, [("Y  =  Xβ  +  Wα  +  ε", "cb")], size=15, sa=8)
P(tf, [("          生物      未知变异", "g")], size=11, sa=14)
P(tf, [("W 只能靠 SVD / PCA 线性估计；负对照基因的信息只在估计时用一次，"
        "无法非线性外推。", "")], size=12.5, sa=0)

tf = tbox(s, ML + CW * 0.53, y, CW * 0.47, 5.4)
P(tf, [("RUVVAE 的做法", "b")], size=16, sa=14)
P(tf, [("把 ", ""), ("W", "cb"), (" 换成 VAE 隐变量，把负对照约束变成一个"
                                  "持续作用的可微损失项。", "")], size=13.5, sa=24)

P(tf, [("log y  =  y_bio  +  Δ_lat  +  Δ_cov", "cb")], size=15, sa=10)
P(tf, [("           生物通道    潜在 UV     已知 UV", "g")], size=11, sa=18)
P(tf, [("y_bio", "cb"), ("    z，64 维 — 细胞层面的真实生物态", "")],
  size=12.5, sa=8)
P(tf, [("Δ_lat", "cb"), ("    w，4 维 + 每家公司一条可学习基线", "")],
  size=12.5, sa=8)
P(tf, [("Δ_cov", "cb"), ("    batch 与测序深度，线性可读", "")], size=12.5, sa=26)
P(tf, [("三通道在 log 尺度上可加 → ", ""), ("y_bio", "cb"),
       (" 直接就是去除 UV 后的干净表达，", ""),
       ("对它做组间比较即得干净 logFC。", "b")], size=13.5, sa=0)

# ================================================================ SLIDE 2
s, y = new_slide(prs, "模型架构 · 三通道可加分解")

tf = tbox(s, ML, y, CW * 0.5, 5.4)
P(tf, [("前向计算", "b")], size=16, sa=14)
P(tf, [("y  →  encoder_z  →  z  (64)", "c")], size=12, sa=6)
P(tf, [("y  →  encoder_w  →  w  (4)", "c")], size=12, sa=16)
P(tf, [("y_bio   = decoder_bio(z) + bias + c_group @ W_group", "c")],
  size=12, sa=6)
P(tf, [("Δ_lat   = sample_emb[batch] + decoder_w(w)", "c")], size=12, sa=6)
P(tf, [("Δ_cov   = c_batch @ W_batch + c_depth @ W_depth", "c")],
  size=12, sa=6)
P(tf, [("y_recon = y_bio + Δ_lat + Δ_cov", "cb")], size=12, sa=24)
P(tf, [("W_group", "cb"), (" 是显式的 2 × 16 428 参数矩阵 —— 组效应线性可读，"
                            "不藏在黑箱里。", "")], size=12.5, sa=30)

P(tf, [("训练配置", "b")], size=16, sa=12)
P(tf, [("AdamW 1e-3 · ReduceLROnPlateau · 150 epochs · batch 256", "g")],
  size=12, sa=6)
P(tf, [("d_bio 64 · k_unk 4 · 重建损失 MSE（另有 ZINB 模式）", "g")],
  size=12, sa=0)

tf = tbox(s, ML + CW * 0.54, y, CW * 0.46, 5.4)
P(tf, [("损失函数", "b"), ("   L = L_recon + KL_z + KL_w + L_NC", "cb"),
       ("  等权", "g")], size=16, sa=14)
P(tf, [("L_recon", "cb"), ("   MSE(y_recon, y)", "c")], size=12, sa=6)
P(tf, [("KL_z", "cb"), ("      z → N(0, I)", "c")], size=12, sa=6)
P(tf, [("KL_w", "cb"), ("      w → N(0, I)，限制 UV 通道容量", "c")],
  size=12, sa=6)
P(tf, [("L_NC", "cb"), ("      MSE(Δ_lat+Δ_cov[nc],  y[nc] − nc_mean)", "ca")],
  size=12, sa=26)

P(tf, [("负对照损失的机制", "b"), ("   全模型的锚", "a")], size=16, sa=12)
P(tf, [("① 假定 500 个负对照基因不受应激影响 → 其全部跨细胞变异都是 UV", "")],
  size=12.5, sa=9)
P(tf, [("② 损失强迫 UV 通道精确吃掉 NC 残差 → UV 幅度被校准，不会过度或不足", "")],
  size=12.5, sa=9)
P(tf, [("③ UV 通道全基因共享 → NC 上学到的模式外推到全部 16 428 个基因", "")],
  size=12.5, sa=0)

add_pic(s, "kl.png", ML, 5.42, 2.5, 1.7)
tf = tbox(s, ML + 2.68, 5.62, CW * 0.5 - 2.68, 1.4)
P(tf, [("KL_w:  0.51 → ", ""), ("0.0084", "ab")], size=12.5, sa=7)
P(tf, [("UV 潜变量被约束得几乎坍缩到先验 —— UV 主要由 "
        "sample_emb 与 W_cov 这些显式通道承担，而非黑箱潜变量。", "g")],
  size=10.5, sa=0)

# ================================================================ SLIDE 3
s, y = new_slide(prs, "诊断 · UV 通道吸收了多少变异")

add_pic(s, "latents.png", ML - 0.1, y + 0.05, CW * 0.58, 3.15)
tf = tbox(s, ML, y + 3.42, CW * 0.56, 1.5)
P(tf, [("左", "b"), ("   生物潜空间 z 按 status 上色 —— 两组大幅重叠、"
                     "重心略有偏移，说明 z 编码的是细胞状态连续谱，"
                     "而非硬分组。", "")], size=11.5, sa=8)
P(tf, [("右", "b"), ("   UV 潜空间 w 按送样公司上色 —— beirui 与 yunzhun "
                     "沿 PC1 明显分开；", ""),
       ("但 seekgene 点数少且与 beirui 大量重叠，分离是部分的。", "ab")],
  size=11.5, sa=0)

tf = tbox(s, ML + CW * 0.60, y, CW * 0.40, 5.6)
P(tf, [("负对照基因找平", "b"), ("   最强证据", "g")], size=15, sa=12)
P(tf, [("跨 donor 标准差   0.0501 → 0.0122   ", "c"), ("0.244×", "ab")],
  size=12, sa=6)
P(tf, [("y_bio[nc] 偏差均值  −0.0002", "c"), ("  |max| 0.0158", "cg")],
  size=12, sa=6)
P(tf, [("重建 MAE          0.0172", "c")], size=12, sa=24)

P(tf, [("重建质量", "b")], size=15, sa=12)
P(tf, [("r²(y_recon)  0.706", "c"), ("     含 UV", "cg")], size=12, sa=6)
P(tf, [("r²(y_bio)    0.307", "c"), ("     去 UV", "cg")], size=12, sa=8)
P(tf, [("落差本身就是 UV 的体量 —— 约 40% 的观测方差是技术性的。", "g")],
  size=11.5, sa=24)

P(tf, [("方差分配", "b"), ("   全基因平均 |logFC|", "g")], size=15, sa=12)
P(tf, [("已知协变量 UV   0.0828", "c")], size=12, sa=6)
P(tf, [("测序深度 UV     0.0822", "c")], size=12, sa=6)
P(tf, [("潜在生物效应    0.0573", "c")], size=12, sa=6)
P(tf, [("显式组效应      0.0255", "c")], size=12, sa=10)
P(tf, [("UV 的量级是生物信号的 2–3 倍。不校正，DEG 基本是在测"
        "「哪家公司测的序」。", "ab")], size=12, sa=0)

# ================================================================ SLIDE 4
s, y = new_slide(
    prs, "核心证据 · Top 20 上调 / Top 20 下调 DEG，按 donor 分组")

tf = tbox(s, ML, y, 3.0, 0.3)
P(tf, [("校正前", "b"), ("   raw log1p", "g")], size=13, sa=0)
add_pic(s, "deg_before.png", ML, y + 0.3, CW, 2.05)

tf = tbox(s, ML, y + 2.48, 4.5, 0.3)
P(tf, [("校正后", "b"), ("   y_bio，UV 已移除", "g")], size=13, sa=0)
add_pic(s, "deg_after.png", ML, y + 2.78, CW, 2.05)

rule(s, ML, 6.42, CW)
tf = tbox(s, ML, 6.58, CW * 0.62, 0.8)
P(tf, [("校正前 7 个 donor 的模式几乎一致，主导变异是 donor 的整体表达水平；"
        "校正后同样这 40 个基因清晰劈成 ", ""),
       ("CSRS(3, CSRES) vs MW(4, CON) 两块", "ab"), ("。", "")],
  size=12, sa=0)
tf = tbox(s, ML + CW * 0.66, 6.58, CW * 0.34, 0.8)
P(tf, [("Welch t-test on logFC_bio · 16 428 基因 · p<0.05: 10 390 · "
        "FDR<0.05: 10 121", "g")], size=10.5, sa=0)

# ================================================================ SLIDE 5
s, y = new_slide(prs, "管家基因验证 · 局限 · 下一步")

tf = tbox(s, ML, y, 4.0, 0.28)
P(tf, [("管家基因 校正前", "b")], size=12, sa=0)
add_pic(s, "hkg_before.png", ML, y + 0.26, CW * 0.56, 1.35)
tf = tbox(s, ML, y + 1.66, 4.0, 0.28)
P(tf, [("管家基因 校正后", "b")], size=12, sa=0)
add_pic(s, "hkg_after.png", ML, y + 1.92, CW * 0.56, 1.35)

tf = tbox(s, ML, y + 3.34, CW * 0.56, 0.9)
P(tf, [("37 个文献管家基因中", "g"), ("仅 4 个", "ab"),
       ("属于负对照集 → 这是真正的 held-out 验证，不是自证。"
        "校正前 MW22B 整体偏低的 donor 梯度，校正后基本消失。", "g")],
  size=11, sa=0)

tf = tbox(s, ML + CW * 0.60, y, CW * 0.40, 5.6)
P(tf, [("需要注意的局限", "ab")], size=15, sa=10)
P(tf, [("FDR<0.05 有 10 121 / 16 428（62%）", "b")], size=11.5, sa=2)
P(tf, [("检验以单细胞为独立样本，存在伪重复；真实独立单位是 7 个 donor，"
        "p 值被极度膨胀。", "g")], size=10.5, sa=8)
P(tf, [("负对照按 |logFC| 最小挑选", "b")], size=11.5, sa=2)
P(tf, [("循环论证风险；选出 Olfr1033、Neurod6 等生物学上不该当 NC 的基因。", "g")],
  size=10.5, sa=8)
P(tf, [("Gm42418 排名第一（logFC 2.43）", "b")], size=11.5, sa=2)
P(tf, [("已知的 rRNA / ambient RNA 污染代理，不应作为生物学结论。", "g")],
  size=10.5, sa=8)
P(tf, [("上下调严重不对称：2 640 上 vs 636 下", "b")], size=11.5, sa=2)
P(tf, [("提示可能仍有未被吸收的全局偏移或 ambient 残留。", "g")],
  size=10.5, sa=8)
P(tf, [("r2_per_gene_mean = −2.5e7", "b")], size=11.5, sa=2)
P(tf, [("被近零方差基因带崩，应以 median 0.081 为准。", "g")], size=10.5, sa=14)

P(tf, [("下一步", "b")], size=15, sa=10)
P(tf, [("改 donor-level pseudobulk 检验", "ab"), ("  ← 最高优先", "g")],
  size=11.5, sa=4)
P(tf, [("负对照改用文献 HKG 池（hkg_only / hybrid）", "")], size=11.5, sa=4)
P(tf, [("过滤 ambient 基因后重跑", "")], size=11.5, sa=4)
P(tf, [("与 MAST / DESeq2-pseudobulk / RUV-III-PRPS 对比", "")], size=11.5, sa=4)
P(tf, [("切换 ZINB 模式直接建模 raw counts", "")], size=11.5, sa=0)

prs.save("RUVVAE_slides.pptx")
print("saved RUVVAE_slides.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst),
      "slides")
